from pptx import Presentation
from pptx.util import Inches


def create_presentation():
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "建設機械プレゼンシナリオ"
    slide.placeholders[1].text = "顧客向け提案資料"

    # Agenda slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "プレゼンの流れ"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. 現場課題の共有"
    tf.add_paragraph().text = "2. 当社製品の特徴と利点"
    tf.add_paragraph().text = "3. 他社との比較"
    tf.add_paragraph().text = "4. 保守サービス"

    # Comparison slide with chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "他社との比較"
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture("comparison.svg", left, top, height=Inches(4))

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "まとめと次のステップ"
    tf = slide.placeholders[1].text_frame
    tf.text = "・ 建設現場での実績"
    tf.add_paragraph().text = "・ デモのご案内"
    tf.add_paragraph().text = "・ お見積り依頼の手順"

    prs.save("machine_sales_pitch.pptx")


if __name__ == "__main__":
    create_presentation()
